import streamlit as st
import openai
import os
import requests
from bs4 import BeautifulSoup
import re
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation

openai.api_key = os.getenv("OPENAI_API_KEY")
client = openai.OpenAI()

if "sow_content" not in st.session_state:
    st.session_state.sow_content = ""

# File extractors
def extract_pdf_text(f): return "".join([p.extract_text() or "" for p in PdfReader(f).pages])[:8000]
def extract_docx_text(f): return "\n".join([p.text for p in Document(f).paragraphs])[:8000]
def extract_excel_text(f): return pd.read_excel(f, sheet_name=None).to_string()[:8000]
def extract_ppt_text(f): return "\n".join(s.text for slide in Presentation(f).slides for s in slide.shapes if hasattr(s, "text"))[:8000]

# Always-on scraping
def fetch_lawinsider():
    try:
        r = requests.get("https://www.lawinsider.com/clause/scope-of-work", timeout=10)
        soup = BeautifulSoup(r.text, "html.parser")
        return [c.get_text(strip=True) for c in soup.select(".clause-body")[:5]]
    except:
        return []

def fetch_text_from_url(url):
    try:
        r = requests.get(url, timeout=10)
        soup = BeautifulSoup(r.text, "html.parser")
        ps = soup.find_all('p')
        return "\n".join(p.get_text(strip=True) for p in ps[:10])
    except Exception as e:
        return f"[Error fetching URL: {e}]"

def fetch_sec_snippets(keyword):
    try:
        search_url = f"https://www.sec.gov/cgi-bin/srch-edgar?text={keyword}"
        r = requests.get(search_url, timeout=10)
        soup = BeautifulSoup(r.text, "html.parser")
        hits = soup.find_all("tr", class_="blueRow")[:2]
        results = []
        base = "https://www.sec.gov"
        for tr in hits:
            link = tr.find("a")
            if link and link.has_attr("href"):
                doc_link = base + link['href']
                doc = requests.get(doc_link, timeout=10)
                doc_soup = BeautifulSoup(doc.text, "html.parser")
                snippet = doc_soup.find("pre") or doc_soup.find("p")
                if snippet:
                    results.append(snippet.get_text(strip=True)[:1500])
        return results
    except:
        return []

def highlight_figures(text):
    return re.sub(r"\$?\d[\d,]*(\.\d+)?", r"[\g<0> – validate independently]", text)

def generate_sow(base_text, desc, examples, sec_snips, role):
    is_vendor = role == "Company is Service Provider"
    example_text = "\n---\n".join(examples + sec_snips)

    prompt = (
        f"You are both a contract lawyer AND a business expert.\n"
        f"Create a detailed Scope of Work (SoW) for the business scenario described below.\n"
        f"Assume a {'Pro-vendor' if is_vendor else 'Pro-client'} stance.\n"
        f"Use 'Company' for client and 'Service Provider' for vendor.\n\n"
        f"Structure:\n"
        f"1. Description\n"
        f"2. Function\n"
        f"3. Price\n"
        f"4. Dependencies\n"
        f"5. Milestones\n"
        f"6. Warranties\n"
        f"7. Service Levels\n"
        f"8. Others\n\n"
        f"---\n"
        f"User Description:\n{desc}\n\n"
        f"---\n"
        f"Base Document Extract:\n{base_text}\n\n"
        f"---\n"
        f"Reference Examples:\n{example_text}\n\n"
        f"Highlight all figures requiring validation.\n"
    )

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "system", "content": "You are a contract lawyer and business expert."},
                  {"role": "user", "content": prompt}],
        temperature=0.5
    )
    return highlight_figures(response.choices[0].message.content)

# Streamlit UI
st.title("AI Scope of Work (SoW) Generator")

file = st.file_uploader("Upload base document (PDF, DOCX, XLSX, PPTX)", type=["pdf", "docx", "xlsx", "xls", "pptx"])
desc = st.text_area("Describe the goods/services and business context")
role = st.radio("Who is Company?", ["Company is Service Provider", "Company is Client"])
kw = st.text_input("Keyword to search SEC filings (e.g. 'drone monitoring')")

url = st.text_input("Paste a URL to extract external SoW-style clauses (optional)")
custom = st.text_area("Paste your own SoW clauses or content here (optional)")

if st.button("Generate Scope of Work"):
    if not (file and desc): st.warning("Please upload a document and provide a description."); st.stop()

    ext = file.name.lower()
    extract_map = {
        "pdf": extract_pdf_text,
        "docx": extract_docx_text,
        "xlsx": extract_excel_text,
        "xls": extract_excel_text,
        "pptx": extract_ppt_text
    }
    base_text = extract_map.get(ext.split('.')[-1], lambda x: "")(file)

    lawinsider = fetch_lawinsider()
    sec_snips = fetch_sec_snippets(kw) if kw.strip() else []
    url_clause = fetch_text_from_url(url) if url.strip() else ""
    custom_clause = custom.strip()

    all_examples = lawinsider
    if url_clause: all_examples.append(url_clause)
    if custom_clause: all_examples.append(custom_clause)

    sow = generate_sow(base_text, desc, all_examples, sec_snips, role)
    st.session_state.sow_content = sow
    st.subheader("Generated Scope of Work")
    st.write(sow)

# Iterative Refinement
st.header("Refine SoW")
if st.session_state.sow_content:
    st.text_area("Current SoW", value=st.session_state.sow_content, height=300, disabled=True)
    feedback = st.text_area("Suggest improvements")
    if st.button("Apply Refinement"):
        if feedback.strip():
            prompt = f"""You are a contract lawyer and business expert. Refine this Scope of Work (SoW) based on user feedback.\n\nCurrent SoW:\n{st.session_state.sow_content}\n\nFeedback:\n{feedback.strip()}\n\nReturn only the refined SoW."""
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are a contract lawyer and business expert."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3
            )
            st.session_state.sow_content = highlight_figures(response.choices[0].message.content)
            st.success("Refinement applied.")
else:
    st.info("Generate a Scope of Work first.")
